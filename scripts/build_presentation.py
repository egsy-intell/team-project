#!/usr/bin/env python3
"""Build the checkpoint slide deck (preso/checkpoint2_deck.md) into a .pptx.

Converts the markdown source into a PowerPoint file styled by
preso/template.pptx (a pandoc --reference-doc), using pypandoc's bundled
pandoc binary (via the `preso` dependency group's pypandoc_binary), so no
system pandoc install is required. The generated .pptx is a build
artifact meant for manual upload to OneDrive - there's no publish step.
"""

import argparse
import os
import subprocess
import sys
from pathlib import Path

try:
    import pypandoc
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print(
        "pypandoc/python-pptx are not installed. Install the "
        "presentation-build dependencies with:\n\n    uv sync --group preso\n",
        file=sys.stderr,
    )
    sys.exit(1)

REPO_ROOT = Path(__file__).resolve().parent.parent
PRESO_DIR = REPO_ROOT / "preso"
SOURCE_MD = PRESO_DIR / "checkpoint2_deck.md"
TEMPLATE_PPTX = PRESO_DIR / "template.pptx"

# Pandoc has no markdown syntax to pick a slide layout by name - it
# auto-selects one per slide from the content's shape (see `pandoc`'s
# manual, "PowerPoint layout choice"). "Content with Caption" is what a
# short intro line followed by a table triggers, and we use it as this
# deck's "dense content" slide type (e.g. References). Table cells don't
# inherit font size from the layout's placeholder style the way plain
# text does, so shrinking them has to happen here, after pandoc has
# actually built the table.
DENSE_CONTENT_LAYOUT = "Content with Caption"
DENSE_TABLE_FONT_SIZE = Pt(14)

# Without an explicit --slide-level, pandoc infers one by scanning for the
# highest (numerically smallest) heading level that's ever immediately
# followed by non-heading content anywhere in the document, and uses *that*
# level - and only that level - to start new slides. Our `#` section
# dividers (Findings So Far, Evaluation Plan & Modeling Proposals, Wrap-Up)
# used to have nothing directly under them but another `#`/`##` heading, so
# pandoc inferred slide-level 2 and every `##` correctly started its own
# slide. The moment a `#` divider gets its own `::: notes :::` block
# directly beneath it (so presenters have a script for that slide too),
# pandoc sees content under a level-1 heading and silently drops the
# inferred slide-level to 1 - collapsing every `##` section back down into
# whichever `#` divider contains it, instead of giving each its own slide.
# Pinning slide-level=2 makes `##` the permanent slide boundary regardless
# of what content ends up under a `#` divider.
SLIDE_LEVEL = 2


def _shrink_dense_tables(pptx_path: Path) -> None:
    prs = Presentation(str(pptx_path))
    for slide in prs.slides:
        if slide.slide_layout.name != DENSE_CONTENT_LAYOUT:
            continue
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = DENSE_TABLE_FONT_SIZE
    prs.save(str(pptx_path))


def _open_in_default_app(path: Path) -> None:
    # Best-effort only: failing to auto-open should never fail the build.
    try:
        if sys.platform == "darwin":
            subprocess.run(["open", str(path)], check=True)
        elif sys.platform == "win32":
            os.startfile(str(path))  # type: ignore[attr-defined]
        else:
            print(f"Auto-open not supported on this platform; open {path} manually.")
    except (OSError, subprocess.SubprocessError) as exc:
        print(f"Could not auto-open {path}: {exc}", file=sys.stderr)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=PRESO_DIR / "dist",
        help="Directory to write the generated .pptx into (default: preso/dist/)",
    )
    parser.add_argument(
        "--open",
        action="store_true",
        help="Open the generated .pptx in the default app afterward (best-effort)",
    )
    parser.add_argument(
        "--template",
        type=Path,
        default=TEMPLATE_PPTX,
        help=(
            "Reference-doc .pptx to style the deck with "
            "(default: preso/template.pptx)"
        ),
    )
    args = parser.parse_args()

    if not SOURCE_MD.exists():
        print(
            f"Missing deck source: {SOURCE_MD.relative_to(REPO_ROOT)}", file=sys.stderr
        )
        return 1
    if not args.template.exists():
        print(
            f"Missing reference template: {args.template}\n"
            "See the README's presentation-pipeline section for how it's built.",
            file=sys.stderr,
        )
        return 1

    output_dir = args.output_dir
    output_dir.mkdir(parents=True, exist_ok=True)
    output_stem = SOURCE_MD.stem
    if args.template != TEMPLATE_PPTX:
        output_stem += f"-{args.template.stem}"
    output_path = output_dir / f"{output_stem}.pptx"

    print(f"Building {SOURCE_MD.relative_to(REPO_ROOT)} -> {output_path}")
    try:
        pypandoc.convert_file(
            str(SOURCE_MD),
            to="pptx",
            outputfile=str(output_path),
            extra_args=[
                "--reference-doc",
                str(args.template),
                "--slide-level",
                str(SLIDE_LEVEL),
            ],
        )
    except RuntimeError as exc:
        print(f"pandoc conversion failed: {exc}", file=sys.stderr)
        return 1

    _shrink_dense_tables(output_path)

    print(f"Wrote {output_path}")

    if args.open:
        _open_in_default_app(output_path)

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
