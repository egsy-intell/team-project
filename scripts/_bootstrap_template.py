"""One-off script that customized ``preso/template.pptx``.

Not part of the ongoing build (``scripts/toolkit.py presentation`` only
needs ``pypandoc`` and reads the committed ``template.pptx`` as-is). Kept
for reference in case the template ever needs to be regenerated (e.g. a
logo refresh) — rerun with ``uv run --with python-pptx python3
scripts/_bootstrap_template.py`` after regenerating the stock template
via ``pandoc -o preso/template.pptx --print-default-data-file
reference.pptx``.

Two customizations:

1. Stamps the team logo onto the layouts pandoc actually uses for this
   deck (large logo on "Title Slide", small mark elsewhere).
2. Widens the "Content with Caption" layout to full slide width. Pandoc
   has no markdown syntax to pick a layout by name — it auto-selects one
   from each slide's content shape (see `pandoc`'s manual, "PowerPoint
   layout choice"), and "text followed by non-text" (e.g. a table) is
   what triggers "Content with Caption". We use that layout as a
   dedicated "dense content" slide type (e.g. References) by writing a
   short intro line + a table in the markdown; the stock version of that
   layout is a narrow caption-plus-picture design, so we widen it to
   full width here to fit a table sensibly.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu, Inches

REPO_ROOT = Path(__file__).resolve().parent.parent
PRESO_DIR = REPO_ROOT / "preso"
TEMPLATE_PPTX = PRESO_DIR / "template.pptx"
LARGE_LOGO = PRESO_DIR / "assets" / "egsy_intell_logo_large.png"
SMALL_LOGO = PRESO_DIR / "assets" / "egsy_intell_logo_small.png"

# Top-right corner mark, sized/positioned to clear each layout's title
# placeholder (confirmed against the stock template's placeholder
# geometry: title boxes start at top >= 0.225in and the widest ones run
# to the right edge, so the mark sits above/right of that text).
LARGE_LOGO_SIZE = Inches(1.1)
LARGE_LOGO_LEFT = Emu(9144000) - LARGE_LOGO_SIZE - Inches(0.35)
LARGE_LOGO_TOP = Inches(0.3)

SMALL_LOGO_SIZE = Inches(0.5)
SMALL_LOGO_LEFT = Emu(9144000) - SMALL_LOGO_SIZE - Inches(0.2)
SMALL_LOGO_TOP = Inches(0.15)

# "Title and Content"'s geometry, reused below so "Content with Caption"
# lines up with every other content slide instead of keeping its stock
# narrow caption-column layout.
FULL_TITLE = (Emu(457200), Emu(205979), Emu(8229600), Emu(857250))
INTRO_LINE = (Emu(457200), Emu(1200151), Emu(8229600), Emu(400050))
WIDE_CONTENT = (Emu(457200), Emu(1650201), Emu(8229600), Emu(2944422))


def _stamp(layout, image_path: Path, left: Emu, top: Emu, size: Emu) -> None:
    # python-pptx's LayoutShapes deliberately omits add_picture (layouts
    # are normally hand-edited in PowerPoint, not scripted) — but the
    # underlying machinery it's built on (SlideShapes.add_picture) works
    # identically for a layout's shape tree, so we replicate it here by
    # reaching into python-pptx internals for this one-off bootstrap.
    shapes = layout.shapes
    image_part, rId = shapes.part.get_or_add_image_part(str(image_path))
    shape_id = shapes._next_shape_id
    scaled_cx, scaled_cy = image_part.scale(size, size)
    name = f"Picture {shape_id - 1}"
    pic = shapes._spTree.add_pic(
        shape_id, name, image_part.desc, rId, left, top, scaled_cx, scaled_cy
    )
    shapes._shape_factory(pic)


def _widen_content_with_caption(layout) -> None:
    geometry_by_idx = {0: FULL_TITLE, 2: INTRO_LINE, 1: WIDE_CONTENT}
    for placeholder in layout.placeholders:
        geometry = geometry_by_idx.get(placeholder.placeholder_format.idx)
        if geometry is None:
            continue
        placeholder.left, placeholder.top, placeholder.width, placeholder.height = (
            geometry
        )


def main() -> int:
    prs = Presentation(str(TEMPLATE_PPTX))
    layouts = {layout.name: layout for layout in prs.slide_masters[0].slide_layouts}

    _widen_content_with_caption(layouts["Content with Caption"])

    _stamp(
        layouts["Title Slide"],
        LARGE_LOGO,
        LARGE_LOGO_LEFT,
        LARGE_LOGO_TOP,
        LARGE_LOGO_SIZE,
    )
    for layout_name in ("Title and Content", "Section Header", "Content with Caption"):
        _stamp(
            layouts[layout_name],
            SMALL_LOGO,
            SMALL_LOGO_LEFT,
            SMALL_LOGO_TOP,
            SMALL_LOGO_SIZE,
        )

    prs.save(str(TEMPLATE_PPTX))
    print(f"Customized {TEMPLATE_PPTX.relative_to(REPO_ROOT)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
