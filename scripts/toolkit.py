#!/usr/bin/env python3
"""Project utility CLI: `presentation` builds the slide deck, `clean-notebook`
patches a marimo notebook export so it prints/PDFs cleanly, `ai-disclosure`
renders a teammate's AI Tool Use Policy disclosure page.

    uv run python scripts/toolkit.py presentation [options]
    uv run python scripts/toolkit.py clean-notebook [options]
    uv run python scripts/toolkit.py ai-disclosure <person> [options]

Run `uv run python scripts/toolkit.py <subcommand> --help` for each
subcommand's options.
"""

import argparse
import json
import os
import re
import subprocess
import sys
import urllib.parse
import urllib.request
from datetime import datetime, timedelta
from pathlib import Path

REPO_ROOT = Path(__file__).resolve().parent.parent

# ---------------------------------------------------------------------------
# `presentation`: build the checkpoint slide deck (preso/checkpoint2_deck.md)
# into a .pptx, styled by preso/template.pptx (a pandoc --reference-doc).
# ---------------------------------------------------------------------------

PRESO_DIR = REPO_ROOT / "preso"
SOURCE_MD = PRESO_DIR / "checkpoint2_deck.md"
TEMPLATE_PPTX = PRESO_DIR / "template.pptx"

# Pandoc has no markdown syntax to pick a slide layout by name - it
# auto-selects one per slide from the content's shape (see `pandoc`'s
# manual, "PowerPoint layout choice"). "Content with Caption" is what a
# short intro line followed by a table triggers, and we use it as this
# deck's "dense content" slide type (e.g. References). Table cells don't
# inherit font size from the layout's placeholder style the way plain
# text does, so shrinking them has to happen here, after pandoc has
# actually built the table.
DENSE_CONTENT_LAYOUT = "Content with Caption"
DENSE_TABLE_FONT_SIZE_PT = 14

# Without an explicit --slide-level, pandoc infers one by scanning for the
# highest (numerically smallest) heading level that's ever immediately
# followed by non-heading content anywhere in the document, and uses *that*
# level - and only that level - to start new slides. Our `#` section
# dividers (Findings So Far, Evaluation Plan & Modeling Proposals, Wrap-Up)
# used to have nothing directly under them but another `#`/`##` heading, so
# pandoc inferred slide-level 2 and every `##` correctly started its own
# slide. The moment a `#` divider gets its own `::: notes :::` block
# directly beneath it (so presenters have a script for that slide too),
# pandoc sees content under a level-1 heading and silently drops the
# inferred slide-level to 1 - collapsing every `##` section back down into
# whichever `#` divider contains it, instead of giving each its own slide.
# Pinning slide-level=2 makes `##` the permanent slide boundary regardless
# of what content ends up under a `#` divider.
SLIDE_LEVEL = 2


def _shrink_dense_tables(pptx_path: Path) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(str(pptx_path))
    for slide in prs.slides:
        if slide.slide_layout.name != DENSE_CONTENT_LAYOUT:
            continue
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            for row in shape.table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(DENSE_TABLE_FONT_SIZE_PT)
    prs.save(str(pptx_path))


def _open_in_default_app(path: Path) -> None:
    # Best-effort only: failing to auto-open should never fail the build.
    try:
        if sys.platform == "darwin":
            subprocess.run(["open", str(path)], check=True)
        elif sys.platform == "win32":
            os.startfile(str(path))  # type: ignore[attr-defined]
        else:
            print(f"Auto-open not supported on this platform; open {path} manually.")
    except (OSError, subprocess.SubprocessError) as exc:
        print(f"Could not auto-open {path}: {exc}", file=sys.stderr)


def cmd_presentation(args: argparse.Namespace) -> int:
    try:
        import pypandoc
    except ImportError:
        print(
            "pypandoc/python-pptx are not installed. Install the "
            "presentation-build dependencies with:\n\n    uv sync --group preso\n",
            file=sys.stderr,
        )
        return 1

    if not SOURCE_MD.exists():
        print(f"Missing deck source: {SOURCE_MD.relative_to(REPO_ROOT)}", file=sys.stderr)
        return 1
    if not args.template.exists():
        print(
            f"Missing reference template: {args.template}\n"
            "See the README's presentation-pipeline section for how it's built.",
            file=sys.stderr,
        )
        return 1

    output_dir = args.output_dir
    output_dir.mkdir(parents=True, exist_ok=True)
    output_stem = SOURCE_MD.stem
    if args.template != TEMPLATE_PPTX:
        output_stem += f"-{args.template.stem}"
    output_path = output_dir / f"{output_stem}.pptx"

    print(f"Building {SOURCE_MD.relative_to(REPO_ROOT)} -> {output_path}")
    try:
        pypandoc.convert_file(
            str(SOURCE_MD),
            to="pptx",
            outputfile=str(output_path),
            extra_args=["--reference-doc", str(args.template), "--slide-level", str(SLIDE_LEVEL)],
        )
    except RuntimeError as exc:
        print(f"pandoc conversion failed: {exc}", file=sys.stderr)
        return 1

    _shrink_dense_tables(output_path)

    print(f"Wrote {output_path}")

    if args.open:
        _open_in_default_app(output_path)

    return 0


def _add_presentation_parser(subparsers: argparse._SubParsersAction) -> None:
    parser = subparsers.add_parser(
        "presentation",
        help="Build the checkpoint slide deck (.pptx) from preso/checkpoint2_deck.md",
        description=cmd_presentation.__doc__ or "",
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=PRESO_DIR / "dist",
        help="Directory to write the generated .pptx into (default: preso/dist/)",
    )
    parser.add_argument(
        "--open",
        action="store_true",
        help="Open the generated .pptx in the default app afterward (best-effort)",
    )
    parser.add_argument(
        "--template",
        type=Path,
        default=TEMPLATE_PPTX,
        help="Reference-doc .pptx to style the deck with (default: preso/template.pptx)",
    )
    parser.set_defaults(func=cmd_presentation)


# ---------------------------------------------------------------------------
# `clean-notebook`: patch a marimo notebook HTML export so it prints/PDFs
# cleanly, instead of the blank pages, clipped tables, and orphaned headings
# marimo's own export produces out of the box.
# ---------------------------------------------------------------------------

NOTEBOOKS_DIR = REPO_ROOT / "notebooks"
DEFAULT_NOTEBOOK = NOTEBOOKS_DIR / "index.py"

# Matches a PEP 723 inline script metadata block, e.g.:
#   # /// script
#   # requires-python = ">=3.14"
#   # ///
# (kept in sync with export_notebooks.py's copy of the same pattern)
_PEP723_HEADER_RE = re.compile(r"^# /// script\n(?:#.*\n)*?# ///\n\n*", re.MULTILINE)

# marimo's HTML export is a single-page app, not a paginated document, and
# none of its layout survives Chrome's print pagination unmodified:
#
#   - The whole notebook sits inside html/body/#root, which marimo pins to
#     one viewport-height with `overflow-y: auto` so the browser can
#     virtually "scroll" it on screen. Printing can't scroll, so anything
#     past the first screenful was clipped or rendered as blank pages.
#   - `mo.vstack()` (and the top-level wrapper around every cell) render as
#     CSS flex containers. Chrome's print engine can't fragment a flex
#     container across a page break, so a tall one gets forced onto a
#     single fresh page as one atomic unit, leaving a large blank gap on
#     the page before it.
#   - The content column is hard-capped at a fixed desktop width (e.g.
#     1110px for marimo's "medium" width), wider than a printed page, so
#     wide tables ran past the margin instead of wrapping.
#   - Headings had no "keep with next" rule, so a heading could land alone
#     at the bottom of a page with its content pushed to the next one.
#
# This CSS (scoped to @media print, so on-screen viewing is untouched) fixes
# all of that. See PRINT_SHADOW_FIX_JS below for the matching fix inside
# shadow DOM, which this stylesheet can't reach.
PRINT_FIX_CSS = """<style id="print-fixes">
@media print {
  @page {
    size: Letter;
    margin: 0.5in;
  }

  /* The app shell pins html/body/#root to one viewport-height, clipped panel
     (overflow-y: auto) so the browser can virtually "scroll" it on screen.
     Printing can't scroll, so everything past the first screenful was either
     blank or cut off. Let the whole chain flow to its natural height. */
  html,
  html > body,
  #root,
  #root > div,
  #root .absolute.inset-0,
  #_r_1_,
  #App {
    height: auto !important;
    max-height: none !important;
    overflow: visible !important;
    position: static !important;
  }

  /* The notebook column is hard-capped at a fixed desktop width (e.g. 1110px
     for "medium"), which is wider than a printed page and pushed content past
     the margins. Let it fill the printable page width instead. */
  :root {
    --content-width-medium: 100% !important;
    --content-width-small: 100% !important;
    --content-width-large: 100% !important;
    --content-width-full: 100% !important;
  }
  [class*="max-w-(--content-width"],
  #App,
  #_r_1_,
  #root > div,
  #root .absolute.inset-0 {
    width: 100% !important;
    max-width: 100% !important;
  }

  /* Every mo.vstack() renders as an inline-styled `display:flex;
     flex-flow:column` box, and Chrome's print engine cannot fragment flex
     containers across a page break: a tall vstack is forced onto a single
     page as one atomic unit, leaving a large blank gap on the page before it
     (or, if it's taller than a page, getting clipped). None of these stacks
     use row layout (marimo emits that separately), so it's safe to drop them
     to normal block flow for print, which fragments cleanly. */
  [style*="flex-flow: column"] {
    display: block !important;
  }

  /* Same problem, but for the one wrapper that lays out the whole notebook
     (all cells, top to bottom) via Tailwind's `flex flex-col` classes rather
     than an inline style: it's tens of thousands of pixels tall, so leaving
     it as a flex container blocks its children from fragmenting across
     pages too. */
  [class~="flex-col"] {
    display: block !important;
  }

  /* Keep headings attached to the content that follows them so a heading
     never lands alone at the bottom of a page (orphaned headers). */
  h1, h2, h3, h4, h5, h6 {
    break-after: avoid !important;
    break-inside: avoid !important;
    page-break-after: avoid !important;
    page-break-inside: avoid !important;
  }

  /* Keep tables, figures, code and math blocks from splitting mid-element
     across a page break. */
  table, marimo-table, marimo-mermaid, marimo-tex, pre, figure, svg {
    break-inside: avoid !important;
    page-break-inside: avoid !important;
  }
  tr {
    break-inside: avoid !important;
    page-break-inside: avoid !important;
  }
  thead {
    display: table-header-group;
  }
  p, li {
    orphans: 3;
    widows: 3;
  }

  /* Floating on-screen chrome (status badges, toast portals) has no meaning
     on paper and otherwise repeats on every printed page. */
  .fixed {
    display: none !important;
  }
}
</style>"""

# marimo custom elements (marimo-table, marimo-mime-renderer for matplotlib
# figures, etc.) render their content inside a shadow root, which the
# document-level stylesheet above can't reach (author stylesheets don't
# cross shadow boundaries). This script walks every shadow root and injects
# the equivalent print-only overrides directly into each one:
#
#   - Some data tables get fixed pixel column widths (sized for the
#     on-screen layout) that add up to wider than a printed page, clipping
#     the right-most columns. table-layout: fixed + width: 100% lets them
#     reflow to fit instead.
#   - matplotlib figures are embedded as full-resolution <img> (routinely
#     750px+ tall). Chrome can't split a raster image across a page break,
#     so an oversized one gets pushed whole onto a fresh page, leaving a
#     large blank gap on the page before it. Capping the height makes it
#     far more likely to fit in whatever room is left.
PRINT_SHADOW_FIX_JS = """<script id="print-shadow-fix">
(function () {
  var PRINT_SHADOW_CSS =
    "@media print {" +
    "table { table-layout: fixed !important; width: 100% !important; } " +
    "th, td { width: auto !important; max-width: none !important; " +
    "white-space: normal !important; overflow-wrap: break-word !important; " +
    "word-break: break-word !important; font-size: 8px !important; " +
    "line-height: 1.25 !important; padding: 2px 4px !important; } " +
    "img { max-height: 6in !important; width: auto !important; " +
    "height: auto !important; max-width: 100% !important; " +
    "object-fit: contain !important; }" +
    "}";
  function injectInto(shadowRoot) {
    if (shadowRoot.getElementById("print-shadow-style")) return;
    var style = document.createElement("style");
    style.id = "print-shadow-style";
    style.textContent = PRINT_SHADOW_CSS;
    shadowRoot.appendChild(style);
  }
  function fixShadowRoots(root) {
    root.querySelectorAll("*").forEach(function (el) {
      if (el.shadowRoot) {
        injectInto(el.shadowRoot);
        fixShadowRoots(el.shadowRoot);
      }
    });
  }
  function run() {
    fixShadowRoots(document);
  }
  window.addEventListener("beforeprint", run);
  window.addEventListener("load", function () {
    setTimeout(run, 500);
  });
})();
</script>"""


class NotebookCleanError(Exception):
    pass


def clean_notebook_html(html: str) -> str:
    """Inject the print-media patch into a marimo notebook HTML export."""
    if 'id="print-fixes"' in html:
        raise NotebookCleanError("input already contains the print-fixes patch")
    if html.count("</head>") != 1:
        raise NotebookCleanError(
            f"expected exactly one </head>, found {html.count('</head>')}"
        )
    if html.count("</body>") != 1:
        raise NotebookCleanError(
            f"expected exactly one </body>, found {html.count('</body>')}"
        )
    html = html.replace("</head>", f"{PRINT_FIX_CSS}\n</head>", 1)
    html = html.replace("</body>", f"{PRINT_SHADOW_FIX_JS}\n</body>", 1)
    return html


def _is_url(ref: str) -> bool:
    return ref.startswith(("http://", "https://"))


def _strip_pep723_header(source: str) -> str:
    # marimo's HTML export embeds the notebook's full source verbatim when
    # code is included (used to reconstruct cells for the editor), so the
    # PEP 723 dependency header would otherwise leak into the exported
    # HTML's embedded source. Only matters when --include-code is used;
    # with the default (no code), marimo omits the source entirely.
    return _PEP723_HEADER_RE.sub("", source, count=1)


def _export_local_notebook(notebook_path: Path, *, include_code: bool) -> str:
    """Run `marimo export html` on a local notebook .py file, return the HTML."""
    print(f"Exporting {notebook_path} (include_code={include_code})")
    source = notebook_path.read_text(encoding="utf-8")
    # Strip the header in place (rather than exporting from a copy
    # elsewhere) so relative sibling imports (e.g. checkpoint_1.py) and
    # data/ lookups keep resolving locally. Always restore afterwards.
    if include_code:
        notebook_path.write_text(_strip_pep723_header(source), encoding="utf-8")
    try:
        # No -o: marimo prints the HTML to stdout when --output is omitted,
        # which we capture directly instead of round-tripping through a
        # temp file. stderr is left to inherit ours, so any warnings the
        # notebook's own code prints during execution still surface live.
        result = subprocess.run(
            [
                sys.executable,
                "-m",
                "marimo",
                "export",
                "html",
                str(notebook_path),
                "--include-code" if include_code else "--no-include-code",
            ],
            stdout=subprocess.PIPE,
            text=True,
            check=False,
        )
    finally:
        if include_code:
            notebook_path.write_text(source, encoding="utf-8")
    if result.returncode != 0:
        raise NotebookCleanError(
            f"marimo export html exited with status {result.returncode} (see output above)"
        )
    return result.stdout


def _resolve_html(input_ref: str, *, include_code: bool) -> str:
    if _is_url(input_ref):
        print(f"Downloading {input_ref}")
        with urllib.request.urlopen(input_ref) as resp:
            return resp.read().decode("utf-8")
    path = Path(input_ref)
    if not path.exists():
        raise NotebookCleanError(f"no such file: {path}")
    if path.suffix == ".py":
        return _export_local_notebook(path, include_code=include_code)
    return path.read_text(encoding="utf-8")


def _default_stem(input_ref: str) -> str:
    # A trailing slash (a directory URL/path) resolves to that directory's
    # index.html.
    if input_ref.endswith("/"):
        return "index"
    name = urllib.parse.urlparse(input_ref).path if _is_url(input_ref) else input_ref
    return Path(name).stem or "index"


def cmd_clean_notebook(args: argparse.Namespace) -> int:
    input_is_local_notebook = not _is_url(args.input) and Path(args.input).suffix == ".py"
    if args.include_code and not input_is_local_notebook:
        print(
            "Warning: --include-code only applies when exporting a local .py "
            "notebook; ignoring (input is already-exported HTML).",
            file=sys.stderr,
        )

    try:
        html = _resolve_html(args.input, include_code=args.include_code)
    except (NotebookCleanError, OSError, urllib.error.URLError) as exc:
        print(f"Failed to read {args.input}: {exc}", file=sys.stderr)
        return 1

    try:
        cleaned = clean_notebook_html(html)
    except NotebookCleanError as exc:
        print(f"Failed to patch {args.input}: {exc}", file=sys.stderr)
        return 1

    name = args.name or f"{_default_stem(args.input)}_clean.html"
    output_dir = args.output_dir or Path.cwd()
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / name

    output_path.write_text(cleaned, encoding="utf-8")
    print(f"Wrote {output_path}")
    return 0


def _add_clean_notebook_parser(subparsers: argparse._SubParsersAction) -> None:
    parser = subparsers.add_parser(
        "clean-notebook",
        help="Export (or read) a marimo notebook HTML export and patch it so it prints/PDFs cleanly",
        description=(
            "Export a marimo notebook to HTML - by default, locally from "
            f"{DEFAULT_NOTEBOOK.relative_to(REPO_ROOT)} - and patch it with a "
            "print-only CSS/JS fix so PDF/print output doesn't have blank pages, "
            "tables clipped past the margin, or orphaned headings. The notebook's "
            "normal on-screen appearance is untouched. INPUT can instead be a "
            "different local notebook .py file, a URL, or an already-exported "
            ".html file (local or downloaded) to patch as-is."
        ),
    )
    parser.add_argument(
        "input",
        nargs="?",
        default=str(DEFAULT_NOTEBOOK),
        help=(
            "Local notebook .py file, local .html file, or URL to a marimo HTML "
            f"export (default: {DEFAULT_NOTEBOOK.relative_to(REPO_ROOT)})"
        ),
    )
    parser.add_argument(
        "--include-code",
        action="store_true",
        help=(
            "Include the notebook's source code in the export (default: excluded, "
            "matching molab). Only applies when INPUT is a local .py notebook."
        ),
    )
    parser.add_argument(
        "--name",
        help="Output filename (default: <input stem>_clean.html)",
    )
    parser.add_argument(
        "--output-dir",
        type=Path,
        help="Directory to write the output into (default: current directory)",
    )
    parser.set_defaults(func=cmd_clean_notebook)


# ---------------------------------------------------------------------------
# `ai-disclosure`: render a teammate's individual AI-use disclosure page
# (docs/ai/<person>.html) and transcript index
# (docs/ai/logs/<person>/index.md) from a small JSON manifest, per the
# course's AI Tool Use Policy (docs/ai/skill/POLICY.md has the full text).
# Deliberately tool-agnostic: nothing here assumes Claude Code specifically
# produced the manifest - see docs/ai/skill/README.md for the runbook
# (including a Claude-Code-specific helper for gathering the manifest's
# thread list from local session logs, and manual guidance for other
# tools). This subcommand only does the deterministic half: manifest (+
# already-written transcripts) in, disclosure page + logs index out.
# ---------------------------------------------------------------------------

REQUIRED_MANIFEST_KEYS = [
    "person", "display_name", "tool", "tool_tier", "how", "why", "threads",
]
REQUIRED_THREAD_KEYS = ["title", "branch", "date", "transcript"]


class ManifestError(Exception):
    pass


def _load_manifest(person: str, ai_logs_dir: Path) -> tuple[dict, Path]:
    manifest_path = ai_logs_dir / person / "_manifest.json"
    if not manifest_path.exists():
        raise ManifestError(
            f"no manifest at {manifest_path} - see "
            "docs/ai/skill/README.md for the manifest format and how to "
            "build one, or docs/ai/skill/manifest.example.json for a "
            "worked example"
        )
    try:
        data = json.loads(manifest_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as exc:
        raise ManifestError(f"{manifest_path}: invalid JSON ({exc})") from exc
    return data, manifest_path


def _validate_manifest(data: dict, manifest_path: Path) -> None:
    missing = [k for k in REQUIRED_MANIFEST_KEYS if k not in data]
    if missing:
        raise ManifestError(f"{manifest_path}: missing required key(s): {', '.join(missing)}")
    if not isinstance(data["threads"], list) or not data["threads"]:
        raise ManifestError(f"{manifest_path}: 'threads' must be a non-empty list")
    uses_phases = bool(data.get("phases"))
    for i, thread in enumerate(data["threads"]):
        missing = [k for k in REQUIRED_THREAD_KEYS if k not in thread]
        if missing:
            raise ManifestError(
                f"{manifest_path}: threads[{i}] missing required key(s): {', '.join(missing)}"
            )
        transcript_path = manifest_path.parent / thread["transcript"]
        if not transcript_path.exists():
            raise ManifestError(
                f"{manifest_path}: threads[{i}]['transcript'] ({thread['transcript']}) "
                f"does not exist in {manifest_path.parent}"
            )
        if "commits" in thread:
            if not isinstance(thread["commits"], list):
                raise ManifestError(f"{manifest_path}: threads[{i}]['commits'] must be a list")
            for j, commit in enumerate(thread["commits"]):
                if not isinstance(commit, dict) or "hash" not in commit or "subject" not in commit:
                    raise ManifestError(
                        f"{manifest_path}: threads[{i}]['commits'][{j}] must be an object "
                        "with 'hash' and 'subject' keys (got "
                        f"{commit!r})"
                    )
        # Phases are all-or-nothing: a thread silently missing from every
        # phase section (because it had no 'phase' key) is worse than a
        # loud error here, since it would just vanish from the rendered
        # page instead of failing the render.
        if "phase" in thread and not uses_phases:
            raise ManifestError(
                f"{manifest_path}: threads[{i}] sets 'phase' but the manifest has no "
                "top-level 'phases' mapping"
            )
        if uses_phases and "phase" not in thread:
            raise ManifestError(
                f"{manifest_path}: threads[{i}] ('{thread.get('title')}') has no 'phase', "
                "but the manifest defines top-level 'phases' - either give every thread a "
                "'phase', or remove 'phases' entirely to render one flat timeline"
            )
        if thread.get("phase") and thread["phase"] not in data.get("phases", {}):
            raise ManifestError(
                f"{manifest_path}: threads[{i}]['phase'] ('{thread['phase']}') is not a key "
                f"in the manifest's top-level 'phases' mapping"
            )


def _git_commits_by_authors(
    author_patterns: list[str], repo_root: Path
) -> list[tuple[str, datetime, str]]:
    """All commits (hash, author-local timestamp, subject) by any of the
    given author name/email substrings, across all branches. Multiple
    --author flags are OR'd by git itself."""
    cmd = ["git", "log", "--all", "--pretty=format:%H%x1f%aI%x1f%s"]
    for pattern in author_patterns:
        cmd += ["--author", pattern]
    result = subprocess.run(cmd, cwd=repo_root, capture_output=True, text=True, check=True)
    commits = []
    for line in result.stdout.splitlines():
        if not line.strip():
            continue
        h, iso, subject = line.split("\x1f", 2)
        commits.append((h, datetime.fromisoformat(iso), subject))
    return commits


def _fill_commits_from_git(data: dict, repo_root: Path, pad_minutes: int = 20) -> None:
    """For any thread with 'start'/'end' timestamps and no explicit
    'commits' list, auto-populate commits from git log by time window.
    Threads that already list 'commits' explicitly are left untouched -
    this only fills gaps, it never overrides a manual list."""
    authors = data.get("git_author_patterns") or [data.get("display_name", "")]
    authors = [a for a in authors if a]
    if not authors:
        return
    all_commits = None
    pad = timedelta(minutes=pad_minutes)
    for thread in data["threads"]:
        if thread.get("commits"):
            continue
        start, end = thread.get("start"), thread.get("end")
        if not start or not end:
            thread.setdefault("commits", [])
            continue
        if all_commits is None:
            all_commits = _git_commits_by_authors(authors, repo_root)
        start_dt = datetime.fromisoformat(start)
        end_dt = datetime.fromisoformat(end)
        matches = [
            (h, ts, subj) for h, ts, subj in all_commits
            if start_dt - pad <= ts <= end_dt + pad
        ]
        thread["commits"] = [
            {"hash": h[:7], "subject": s, "time": ts.strftime("%H:%M")}
            for h, ts, s in matches
        ]


TIMELINE_CSS = """
  @font-face {
    font-family: "ui-mono";
    src: local("SF Mono"), local("Cascadia Code"), local("Consolas");
  }

  :root {
    --bg: #f4f7f6;
    --surface: #ffffff;
    --surface-2: #eaf0ee;
    --ink: #12201c;
    --ink-soft: #4f625c;
    --ink-faint: #7c8c86;
    --line: #d7e2de;
    --accent: #1c7268;
    --accent-soft: #e3f0ed;
    --accent-ink: #0f4a44;
    color-scheme: light;
  }

  @media (prefers-color-scheme: dark) {
    :root {
      --bg: #0c1412; --surface: #121d1a; --surface-2: #172521;
      --ink: #e7efec; --ink-soft: #a7bab3; --ink-faint: #6d827b;
      --line: #22332e; --accent: #45b0a3; --accent-soft: #16302c;
      --accent-ink: #8fdccf; color-scheme: dark;
    }
  }
  :root[data-theme="dark"] {
    --bg: #0c1412; --surface: #121d1a; --surface-2: #172521;
    --ink: #e7efec; --ink-soft: #a7bab3; --ink-faint: #6d827b;
    --line: #22332e; --accent: #45b0a3; --accent-soft: #16302c;
    --accent-ink: #8fdccf; color-scheme: dark;
  }
  :root[data-theme="light"] {
    --bg: #f4f7f6; --surface: #ffffff; --surface-2: #eaf0ee;
    --ink: #12201c; --ink-soft: #4f625c; --ink-faint: #7c8c86;
    --line: #d7e2de; --accent: #1c7268; --accent-soft: #e3f0ed;
    --accent-ink: #0f4a44; color-scheme: light;
  }

  * { box-sizing: border-box; }
  body {
    margin: 0; background: var(--bg); color: var(--ink);
    font-family: -apple-system, "Segoe UI", system-ui, "Helvetica Neue", Arial, sans-serif;
    line-height: 1.5; -webkit-font-smoothing: antialiased;
  }
  .wrap { max-width: 800px; margin: 0 auto; padding: 4rem 1.5rem 6rem; overflow-x: hidden; }
  code { font-family: ui-mono, "SF Mono", "Cascadia Code", Consolas, monospace; font-size: 0.92em; }
  a { color: var(--accent-ink); text-decoration-color: var(--line); }
  a:hover { text-decoration-color: currentColor; }
  a:focus-visible, button:focus-visible { outline: 2px solid var(--accent); outline-offset: 2px; }
  .eyebrow {
    font-size: 0.72rem; font-weight: 700; letter-spacing: 0.12em;
    text-transform: uppercase; color: var(--accent-ink); margin: 0 0 0.75rem;
  }
  h1 { font-size: clamp(1.9rem, 4vw, 2.5rem); line-height: 1.15; letter-spacing: -0.01em; margin: 0 0 0.9rem; text-wrap: balance; }
  .lede { font-size: 1.05rem; color: var(--ink-soft); max-width: 62ch; margin: 0 0 2rem; }
  .lede strong { color: var(--ink); font-weight: 600; }
  .stats {
    display: grid; grid-template-columns: repeat(auto-fit, minmax(120px, 1fr));
    gap: 1px; background: var(--line); border: 1px solid var(--line);
    border-radius: 10px; overflow: hidden; margin-bottom: 1.75rem;
  }
  .stat { background: var(--surface); padding: 0.9rem 1rem; }
  .stat .n {
    font-family: ui-mono, monospace; font-variant-numeric: tabular-nums;
    font-size: 1.35rem; font-weight: 600; color: var(--accent-ink); display: block;
  }
  .stat .l { font-size: 0.72rem; color: var(--ink-faint); text-transform: uppercase; letter-spacing: 0.06em; }
  .note {
    font-size: 0.88rem; color: var(--ink-faint); border-left: 2px solid var(--line);
    padding-left: 0.9rem; margin: 0 0 2rem;
  }
  .disclosure {
    display: grid; grid-template-columns: repeat(2, 1fr); gap: 1px;
    background: var(--line); border: 1px solid var(--line); border-radius: 10px;
    overflow: hidden; margin: 0 0 1.75rem;
  }
  .disclosure-item { background: var(--surface); padding: 1rem 1.15rem; }
  .disclosure-item .dl-label {
    display: block; font-size: 0.68rem; font-weight: 700; letter-spacing: 0.08em;
    text-transform: uppercase; color: var(--accent-ink); margin-bottom: 0.4rem;
  }
  .disclosure-item p { margin: 0; font-size: 0.86rem; color: var(--ink-soft); }
  @media (max-width: 560px) { .disclosure { grid-template-columns: 1fr; } }
  .legend {
    display: flex; flex-wrap: wrap; gap: 1.25rem; margin: 0 0 3rem;
    font-size: 0.78rem; color: var(--ink-faint);
  }
  .legend span { display: inline-flex; align-items: center; gap: 0.4rem; }
  .legend .dot {
    width: 9px; height: 9px; border-radius: 50%; background: var(--accent);
    display: inline-block; flex: none; opacity: 0.55;
  }
  .legend .dot.thread { width: 11px; height: 11px; opacity: 1; }
  .legend .dot.milestone {
    width: 13px; height: 13px; background: var(--accent-soft);
    border: 1.5px solid var(--accent); opacity: 1;
  }
  .toc {
    display: flex; flex-wrap: wrap; gap: 0.4rem; margin: 0 0 3.5rem;
    padding: 0; list-style: none;
  }
  .toc a {
    display: block; font-size: 0.82rem; padding: 0.35rem 0.7rem;
    border: 1px solid var(--line); border-radius: 999px; color: var(--ink-soft);
    text-decoration: none; white-space: nowrap;
  }
  .toc a:hover { border-color: var(--accent); color: var(--accent-ink); }
  section.phase { margin-bottom: 2.75rem; }
  .tl-row { display: flex; align-items: flex-start; }
  .tl-time {
    flex: 0 0 4.75rem; width: 4.75rem; text-align: right; padding-right: 0.9rem;
    padding-top: 0.2rem; font-family: ui-mono, monospace; font-variant-numeric: tabular-nums; line-height: 1.4;
  }
  .tl-time .d { display: block; font-size: 0.72rem; color: var(--ink-soft); font-weight: 600; }
  .tl-time .t { display: block; font-size: 0.66rem; color: var(--ink-faint); }
  .tl-rail { flex: 0 0 1.6rem; width: 1.6rem; position: relative; align-self: stretch; }
  .tl-rail::before {
    content: ""; position: absolute; left: 50%; top: 0; bottom: 0; width: 1px;
    background: var(--line); transform: translateX(-50%);
  }
  .tl-rail .dot {
    position: absolute; left: 50%; top: 0.55rem; width: 7px; height: 7px; border-radius: 50%;
    background: var(--accent); opacity: 0.55; transform: translate(-50%, -50%); box-shadow: 0 0 0 3px var(--bg);
  }
  .tl-row.thread .tl-rail .dot { width: 10px; height: 10px; opacity: 1; }
  .tl-row.milestone .tl-rail::before { top: 0.7rem; }
  .tl-row.milestone .tl-rail .dot {
    width: 1.55rem; height: 1.55rem; border-radius: 50%; background: var(--accent-soft);
    border: 1.5px solid var(--accent); opacity: 1; top: 0.7rem; display: flex; align-items: center;
    justify-content: center; font-family: ui-mono, monospace; font-size: 0.65rem; font-weight: 700;
    color: var(--accent-ink); box-shadow: 0 0 0 3px var(--bg);
  }
  .tl-content { flex: 1 1 auto; min-width: 0; padding-bottom: 1.2rem; }
  .tl-row.milestone .tl-content { padding-top: 0.2rem; padding-bottom: 0.9rem; }
  .tl-row.thread .tl-content { padding-bottom: 0.7rem; }
  .tl-row.milestone h2 { font-size: 1.15rem; margin: 0 0 0.35rem; letter-spacing: -0.005em; scroll-margin-top: 1.5rem; }
  .phase-blurb { color: var(--ink-soft); font-size: 0.94rem; margin: 0; max-width: 58ch; }
  .tl-row.thread h3 { font-size: 1rem; font-weight: 700; margin: 0 0 0.3rem; letter-spacing: -0.003em; }
  .thread-quote {
    margin: 0 0 0.4rem; padding: 0.15rem 0 0.15rem 0.75rem; border-left: 2px solid var(--accent);
    color: var(--ink-soft); font-size: 0.87rem; font-style: italic; max-width: 56ch;
  }
  .thread-meta { font-family: ui-mono, monospace; font-size: 0.68rem; color: var(--ink-faint); font-variant-numeric: tabular-nums; }
  ol.entries { list-style: none; margin: 0 0 0.3rem; padding: 0; }
  .entry-head { display: flex; align-items: baseline; justify-content: space-between; gap: 0.75rem; }
  .entry-head h4 { font-size: 0.93rem; font-weight: 600; margin: 0 0 0.28rem; letter-spacing: -0.003em; color: var(--ink); }
  .entry-head .hash { font-family: ui-mono, monospace; font-size: 0.66rem; color: var(--ink-faint); flex: 0 0 auto; padding-top: 0.12rem; }
  .tl-content p.entry-p { margin: 0 0 0.4rem; color: var(--ink-soft); font-size: 0.87rem; max-width: 55ch; }
  .diffstat { font-family: ui-mono, monospace; font-size: 0.68rem; color: var(--ink-faint); font-variant-numeric: tabular-nums; }
  .diffstat .plus { color: var(--accent-ink); }
  @media (max-width: 620px) {
    .wrap { padding: 2.5rem 1rem 4rem; }
    .tl-time { flex-basis: 3.4rem; width: 3.4rem; padding-right: 0.5rem; }
    .stats { grid-template-columns: repeat(2, 1fr); }
    .entry-head { flex-direction: column; gap: 0.1rem; }
    .thread-meta { line-height: 1.5; }
  }
  @media (max-width: 400px) {
    .tl-time { flex-basis: 2.5rem; width: 2.5rem; padding-right: 0.4rem; }
    .tl-rail { flex-basis: 1.2rem; width: 1.2rem; }
    .stats { grid-template-columns: 1fr 1fr; }
  }
  footer.page { border-top: 1px solid var(--line); padding-top: 1.5rem; color: var(--ink-faint); font-size: 0.82rem; }
  footer.page a { color: inherit; }
"""


def _esc(text: str) -> str:
    return (
        text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        .replace('"', "&quot;")
    )


def _split_day_time(date_str: str) -> tuple[str, str | None]:
    # Display dates are either a bare day ("Aug 4") or "day, time"
    # ("Aug 4, 22:56") - split so the time can render in its own
    # smaller/fainter span, matching how commit timestamps are shown.
    if "," in date_str:
        day, _, time = date_str.partition(",")
        return day.strip(), time.strip()
    return date_str, None


def _render_time_cell(date_str: str) -> str:
    day, time = _split_day_time(date_str)
    time_html = f'<span class="t">{_esc(time)}</span>' if time else ""
    return f'<span class="d">{_esc(day)}</span>{time_html}'


def _render_thread_row(thread: dict, logs_dir_name: str) -> str:
    title = _esc(thread["title"])
    branch = _esc(thread["branch"])
    thread_day, _ = _split_day_time(thread["date"])
    quote_html = ""
    if thread.get("quote"):
        quote_html = f'\n        <p class="thread-quote">&ldquo;{_esc(thread["quote"])}&rdquo;</p>'
    meta_bits = [f"branch {branch}"]
    if thread.get("prompts") is not None and thread.get("responses") is not None:
        meta_bits.append(f"{thread['prompts']} prompts, {thread['responses']} responses")
    meta_bits.append(
        f'<a href="logs/{logs_dir_name}/{_esc(thread["transcript"])}">full transcript &rarr;</a>'
    )
    meta = " &middot; ".join(meta_bits)

    entries = ""
    for commit in thread.get("commits", []):
        time_html = ""
        if commit.get("time"):
            time_html = (
                f'<span class="d">{_esc(thread_day)}</span>'
                f'<span class="t">{_esc(commit["time"])}</span>'
            )
        summary_html = ""
        if commit.get("summary"):
            summary_html = f'\n          <p class="entry-p">{_esc(commit["summary"])}</p>'
        diffstat_html = ""
        if commit.get("diffstat"):
            diffstat_html = f'\n          <div class="diffstat">{commit["diffstat"]}</div>'
        entries += f"""      <li class="tl-row entry">
        <div class="tl-time">{time_html}</div>
        <div class="tl-rail"><span class="dot"></span></div>
        <div class="tl-content">
          <div class="entry-head"><h4>{_esc(commit["subject"])}</h4><span class="hash">{_esc(commit["hash"])}</span></div>{summary_html}{diffstat_html}
        </div>
      </li>
"""
    entries_block = f'\n    <ol class="entries">\n{entries}    </ol>' if entries else ""

    return f"""    <div class="tl-row thread">
      <div class="tl-time">{_render_time_cell(thread["date"])}</div>
      <div class="tl-rail"><span class="dot"></span></div>
      <div class="tl-content">
        <h3>{title}</h3>{quote_html}
        <p class="thread-meta">{meta}</p>
      </div>
    </div>{entries_block}
"""


def render_timeline_html(data: dict) -> str:
    person = data["person"]
    display_name = _esc(data["display_name"])
    tool = _esc(data["tool"])
    tool_tier = _esc(data["tool_tier"])
    how = _esc(data["how"])
    why = _esc(data["why"])

    n_threads = len(data["threads"])
    n_commits = sum(len(t.get("commits", [])) for t in data["threads"])

    stats = [f'<div class="stat"><span class="n">{n_threads}</span><span class="l">Threads</span></div>']
    stats.append(f'<div class="stat"><span class="n">{n_commits}</span><span class="l">Commits with {tool}</span></div>')
    if data.get("total_commits"):
        pct = round(100 * n_commits / data["total_commits"])
        stats.append(f'<div class="stat"><span class="n">{pct}%</span><span class="l">Of {display_name.split()[0]}\'s {data["total_commits"]} commits</span></div>')
    if data.get("date_range") and data.get("days"):
        stats.append(
            f'<div class="stat"><span class="n">{data["days"]}</span>'
            f'<span class="l">Days, {_esc(data["date_range"])}</span></div>'
        )
    elif data.get("date_range"):
        stats.append(f'<div class="stat"><span class="n">{_esc(data["date_range"])}</span><span class="l">Span</span></div>')

    toc_html = ""
    if data.get("phases"):
        body = ""
        seen_phases = []
        for t in data["threads"]:
            phase_key = t.get("phase")
            if phase_key and phase_key not in seen_phases:
                seen_phases.append(phase_key)
        toc_items = "".join(
            f'    <li><a href="#p{i}">{i}. {_esc(data["phases"][phase_key]["title"])}</a></li>\n'
            for i, phase_key in enumerate(seen_phases, start=1)
        )
        toc_html = f'  <ul class="toc">\n{toc_items}  </ul>\n\n'
        for i, phase_key in enumerate(seen_phases, start=1):
            phase = data["phases"][phase_key]
            body += f"""  <section class="phase">
    <div class="tl-row milestone" id="p{i}">
      <div class="tl-time"></div>
      <div class="tl-rail"><span class="dot">{i:02d}</span></div>
      <div class="tl-content">
        <h2>{_esc(phase["title"])}</h2>
        <p class="phase-blurb">{_esc(phase.get("blurb", ""))}</p>
      </div>
    </div>
"""
            for t in data["threads"]:
                if t.get("phase") == phase_key:
                    body += _render_thread_row(t, person)
            body += "  </section>\n\n"
    else:
        body = '  <section class="phase">\n'
        for t in data["threads"]:
            body += _render_thread_row(t, person)
        body += "  </section>\n\n"

    policy_link = data.get(
        "policy_url",
        "https://purdue.brightspace.com/d2l/le/content/1565125/viewContent/21824036/View",
    )

    lede = data.get("lede") or (
        f"This page is {display_name}'s individual pair-programming disclosure under "
        f'the course\'s <a href="{_esc(policy_link)}">AI Tool Use Policy</a>, which asks '
        "for exactly which tools were used and their tier, the history of each "
        "exchange, and how and why each tool was used."
    )
    methodology_note = data.get("methodology_note") or (
        f"Generated from a manifest at "
        f"<code>docs/ai/logs/{person}/_manifest.json</code> via "
        f"<code>uv run python scripts/toolkit.py ai-disclosure {person}</code> - see "
        f'<a href="skill/README.md">docs/ai/skill/</a> for how that manifest and its '
        f'transcripts were put together, and <a href="skill/POLICY.md">the policy '
        f"text</a> this disclosure responds to."
    )
    footer_note = data.get("footer_note") or (
        "This page is generated, not hand-authored - see "
        "<code>docs/ai/skill/README.md</code> to regenerate it after adding "
        "threads or commits."
    )
    legend_html = ""
    if data.get("phases"):
        legend_html = """
  <div class="legend">
    <span><span class="dot milestone"></span> phase</span>
    <span><span class="dot thread"></span> conversation thread</span>
    <span><span class="dot"></span> commit</span>
  </div>
"""

    return f"""<title>Pairing Log &mdash; {display_name}</title>
<meta name="viewport" content="width=device-width, initial-scale=1">
<meta name="description" content="AI use disclosure for {display_name}, per the course AI Tool Use Policy.">
<style>{TIMELINE_CSS}</style>

<div class="wrap">

  <p class="eyebrow">AI Use Disclosure</p>
  <h1>{data.get("h1") or f"Every thread and commit where {tool} paired with {display_name.split()[0]} on this project"}</h1>
  <p class="lede">
    {lede}
  </p>

  <div class="disclosure">
    <div class="disclosure-item">
      <span class="dl-label">Tool &amp; tier</span>
      <p>{tool} &mdash; {tool_tier}.</p>
    </div>
    <div class="disclosure-item">
      <span class="dl-label">How</span>
      <p>{how}</p>
    </div>
    <div class="disclosure-item">
      <span class="dl-label">Why</span>
      <p>{why}</p>
    </div>
    <div class="disclosure-item">
      <span class="dl-label">History of the exchange</span>
      <p>Every thread below links its full prompt/response transcript; all
      {n_threads} are also collected in
      <a href="logs/{person}/index.md"><code>docs/ai/logs/{person}/</code></a>.</p>
    </div>
  </div>

  <div class="stats">
    {"".join(stats)}
  </div>

  <p class="note">
    <strong>Methodology.</strong> {methodology_note}
  </p>
{legend_html}
{toc_html}{body}  <footer class="page">
    <p>
      {footer_note}
    </p>
  </footer>

</div>
"""


def render_logs_index_md(data: dict) -> str:
    person = data["person"]
    display_name = data["display_name"]
    lines = [
        f"# Pair-programming transcripts &mdash; {display_name}".replace("&mdash;", "—"),
        "",
        (
            f"Full prompt/response history for every {data['tool']} thread behind "
            f"[{display_name.split()[0]}'s pairing-log timeline](../../{person}.html), "
            "one file per thread. Text-only: prompts and responses, no "
            "tool-call/tool-result mechanics."
        ),
        "",
    ]
    for i, t in enumerate(data["threads"], start=1):
        lines.append(f"{i}. [{t['title']}]({t['transcript']})")
    lines.append("")
    lines.append(
        "Other teammates' disclosures live in their own sibling directory "
        "under `docs/ai/logs/`."
    )
    return "\n".join(lines) + "\n"


def render_readme_snippet(data: dict) -> str:
    person = data["person"]
    n_threads = len(data["threads"])
    n_commits = sum(len(t.get("commits", [])) for t in data["threads"])
    thread_word = "thread" if n_threads == 1 else "threads"
    commit_word = "commit" if n_commits == 1 else "commits"
    return (
        f"- **{data['display_name']}** &mdash; {data['tool']} ({data['tool_tier']}). "
        f"History of exchange: [`docs/ai/{person}.html`]"
        f"(https://egsy-intell.github.io/team-project/ai/{person}.html) "
        f"({n_threads} {thread_word}, {n_commits} {commit_word}), full transcripts in "
        f"[`docs/ai/logs/{person}/`](docs/ai/logs/{person}/index.md). "
        f"How: {data['how']} Why: {data['why']}\n"
    ).replace("&mdash;", "—")


def cmd_ai_disclosure(args: argparse.Namespace) -> int:
    repo_root = args.repo_root
    ai_dir = repo_root / "docs" / "ai"
    ai_logs_dir = ai_dir / "logs"

    try:
        data, manifest_path = _load_manifest(args.person, ai_logs_dir)
        _validate_manifest(data, manifest_path)
        if not args.skip_git:
            _fill_commits_from_git(data, repo_root=repo_root)
    except ManifestError as exc:
        print(f"Error: {exc}", file=sys.stderr)
        return 1

    person_dir = ai_logs_dir / args.person
    html_path = ai_dir / f"{args.person}.html"
    index_path = person_dir / "index.md"
    snippet_path = person_dir / "_readme_snippet.md"

    html_path.write_text(render_timeline_html(data), encoding="utf-8")
    index_path.write_text(render_logs_index_md(data), encoding="utf-8")
    snippet_path.write_text(render_readme_snippet(data), encoding="utf-8")

    print(f"Wrote {html_path}")
    print(f"Wrote {index_path}")
    print(f"Wrote {snippet_path} (paste into README.md's Pair-programming sessions list)")
    return 0


def _add_ai_disclosure_parser(subparsers: argparse._SubParsersAction) -> None:
    parser = subparsers.add_parser(
        "ai-disclosure",
        help="Render a teammate's AI-use disclosure page + logs index from docs/ai/logs/<person>/_manifest.json",
        description=(
            "Render docs/ai/<person>.html and docs/ai/logs/<person>/index.md "
            "from docs/ai/logs/<person>/_manifest.json and that directory's "
            "transcript files. Tool-agnostic - see docs/ai/skill/README.md "
            "for the manifest format and how to build one with any AI "
            "coding assistant, not just Claude Code."
        ),
    )
    parser.add_argument(
        "person",
        help="Directory name under docs/ai/logs/ holding _manifest.json and transcripts (e.g. a Purdue ID)",
    )
    parser.add_argument(
        "--repo-root",
        type=Path,
        default=REPO_ROOT,
        help=argparse.SUPPRESS,  # mainly for tests; real usage always wants this repo
    )
    parser.add_argument(
        "--skip-git",
        action="store_true",
        help="Don't auto-fill commits from git log by time window; use only what's already listed in the manifest",
    )
    parser.set_defaults(func=cmd_ai_disclosure)


# ---------------------------------------------------------------------------


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    subparsers = parser.add_subparsers(dest="command", required=True)
    _add_presentation_parser(subparsers)
    _add_clean_notebook_parser(subparsers)
    _add_ai_disclosure_parser(subparsers)

    args = parser.parse_args()
    return args.func(args)


if __name__ == "__main__":
    raise SystemExit(main())
